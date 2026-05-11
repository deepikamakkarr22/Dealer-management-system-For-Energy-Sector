"""
Insert a 'Scenario' slide at index 5 — before 'SQL Normalized Schema'.
Shows the e-commerce entities (Customer, Order, Order Item),
their relationships, and what the audience will see modeled next.
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import copy, os

# ── colours ───────────────────────────────────────────────────────────────────
def rgb(h): return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))

BG      = "001E2B"
GREEN   = "00ED64"
WHITE   = "FFFFFF"
GRAY    = "B8C4C2"
CARD    = "023430"
DARK2   = "012A3A"
BLUE    = "4A9EFF"
ORANGE  = "F97316"
PURPLE  = "A855F7"

F_HEAD  = "Lexend Deca"
F_CODE  = "Source Code Pro"
IN = 914400

def emu(i): return int(i * IN)

# ── low-level shape helpers ───────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, border=None, bpt=1.5):
    sp = slide.shapes.add_shape(1, emu(x), emu(y), emu(w), emu(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb(fill)
    if border: sp.line.color.rgb = rgb(border); sp.line.width = Pt(bpt)
    else:       sp.line.fill.background()
    return sp

def txt(slide, text, x, y, w, h,
        size=14, bold=False, color=WHITE, font=F_HEAD,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tb.word_wrap = wrap
    tf = tb.text_frame; tf.word_wrap = wrap
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run()
    r.text = text; r.font.name = font; r.font.size = Pt(size)
    r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = rgb(color)
    return tb

def multiline(slide, lines, x, y, w, h,
              size=13, bold=False, color=WHITE, font=F_HEAD,
              align=PP_ALIGN.LEFT, line_gap=2):
    tb = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(line_gap)
        r = p.add_run()
        r.text = line; r.font.name = font; r.font.size = Pt(size)
        r.font.bold = bold; r.font.color.rgb = rgb(color)
    return tb

# ═════════════════════════════════════════════════════════════════════════════
# Build the slide in a temp presentation
# ═════════════════════════════════════════════════════════════════════════════
tmp = Presentation()
tmp.slide_width  = Emu(emu(13.33))
tmp.slide_height = Emu(emu(7.5))

s = tmp.slides.add_slide(tmp.slide_layouts[6])   # blank

# Background
s.background.fill.solid()
s.background.fill.fore_color.rgb = rgb(BG)

# Bottom green bar
rect(s, 0, 7.28, 13.33, 0.22, GREEN)

# ── Heading ───────────────────────────────────────────────────────────────────
txt(s, "The Scenario — E-Commerce Order Management",
    0.62, 0.2, 12.0, 0.65,
    size=30, bold=True, color=WHITE)
rect(s, 0.62, 0.88, 1.4, 0.06, GREEN)   # underline

txt(s, "Every concept from here onward — SQL normalization, MongoDB modeling, "
       "queries, indexes, transactions — uses these three entities:",
    0.62, 0.98, 12.0, 0.48,
    size=15, color=GRAY)

# ═════════════════════════════════════════════════════════════════════════════
# ENTITY BOXES  (3 cards)
# ═════════════════════════════════════════════════════════════════════════════
EW, EH = 3.5, 3.55

def entity_box(slide, x, label, fields, accent, icon):
    """Draw one entity box."""
    y = 1.58
    # card body
    rect(slide, x, y, EW, EH, CARD, accent, 2.0)
    # header
    rect(slide, x, y, EW, 0.65, accent)
    # icon + name
    txt(slide, icon + "  " + label, x+0.14, y+0.06, EW-0.28, 0.55,
        size=18, bold=True, color=BG, font=F_HEAD)
    # fields
    for i, (fname, ftype) in enumerate(fields):
        fy = y + 0.75 + i * 0.52
        # alternate row shade
        if i % 2 == 0:
            rect(slide, x+0.08, fy-0.04, EW-0.16, 0.46, DARK2)
        txt(slide, fname, x+0.18, fy, 2.1, 0.44,
            size=12, color=WHITE, font=F_CODE)
        txt(slide, ftype, x+2.35, fy, EW-2.5, 0.44,
            size=11, color=GRAY, font=F_CODE, align=PP_ALIGN.RIGHT)

entity_box(s, 0.62, "Customer",
           [("id",         "PK"),
            ("name",       "VARCHAR"),
            ("email",      "VARCHAR"),
            ("region",     "VARCHAR"),
            ("tier",       "ENUM"),],
           GREEN, "👤")

entity_box(s, 4.92, "Order",
           [("id",          "PK"),
            ("customer_id", "FK → Customer"),
            ("status",      "VARCHAR"),
            ("amount",      "DECIMAL"),
            ("created_at",  "TIMESTAMP"),],
           BLUE, "📦")

entity_box(s, 9.20, "Order Item",
           [("id",          "PK"),
            ("order_id",    "FK → Order"),
            ("sku",         "VARCHAR"),
            ("qty",         "INT"),
            ("unit_price",  "DECIMAL"),],
           ORANGE, "🛒")

# ── Relationship arrows ───────────────────────────────────────────────────────
# Customer → Order  (1:N)
rect(s, 4.12, 3.32, 0.8, 0.06, GREEN)       # line
rect(s, 4.88, 3.22, 0.06, 0.26, GREEN)      # arrowhead
txt(s, "places", 4.12, 2.88, 0.8, 0.35,
    size=12, color=GREEN, italic=True, align=PP_ALIGN.CENTER)
txt(s, "1 : N",  4.1, 3.46, 0.85, 0.3,
    size=13, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# Order → OrderItem  (1:N)
rect(s, 8.42, 3.32, 0.78, 0.06, BLUE)
rect(s, 9.16, 3.22, 0.06, 0.26, BLUE)
txt(s, "contains", 8.38, 2.88, 0.86, 0.35,
    size=12, color=BLUE, italic=True, align=PP_ALIGN.CENTER)
txt(s, "1 : N",    8.38, 3.46, 0.86, 0.3,
    size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# ═════════════════════════════════════════════════════════════════════════════
# WHAT WE'LL COVER section  (bottom 3 cards)
# ═════════════════════════════════════════════════════════════════════════════
txt(s, "What you'll see modeled with these entities:",
    0.62, 5.26, 12.0, 0.4,
    size=14, bold=True, color=WHITE)

cards = [
    { "n":"01", "col": GREEN,  "title":"SQL  —  3NF Normalized",
      "desc":"3 separate tables · Foreign keys · JOIN to read one order" },
    { "n":"02", "col": BLUE,   "title":"MongoDB  —  Embedded Documents",
      "desc":"1 document · Order item array embedded · Single read, no joins" },
    { "n":"03", "col": ORANGE, "title":"Queries · Indexes · Transactions",
      "desc":"Same entities, same data — side-by-side SQL vs MongoDB examples" },
]

for i, c in enumerate(cards):
    cx = 0.62 + i * 4.24
    rect(s, cx, 5.72, 4.1, 1.3, CARD, c["col"], 1.5)
    rect(s, cx, 5.72, 4.1, 0.12, c["col"])
    txt(s, c["n"] + "  " + c["title"],
        cx+0.18, 5.8, 3.8, 0.5,
        size=14, bold=True, color=WHITE)
    txt(s, c["desc"],
        cx+0.18, 6.32, 3.8, 0.6,
        size=12, color=GRAY, wrap=True)

# ─────────────────────────────────────────────────────────────────────────────
tmp.save("_scenario_tmp.pptx")
print("Temp slide built")

# ═════════════════════════════════════════════════════════════════════════════
# Merge into the main PPTX at index 5
# ═════════════════════════════════════════════════════════════════════════════
main = Presentation("MongoDB_for_Relational_Developers.pptx")
src  = Presentation("_scenario_tmp.pptx")

def inject_slide(dest_prs, src_prs, src_idx, insert_at):
    """Copy src slide and insert it at insert_at position in dest_prs."""
    template = src_prs.slides[src_idx]

    # Add a blank slide at the END of dest_prs
    blank = dest_prs.slide_layouts[0]
    new_slide = dest_prs.slides.add_slide(blank)

    # Strip placeholder shapes the layout added
    sp_tree = new_slide.shapes._spTree
    for sp in list(sp_tree.findall(qn('p:sp'))):
        sp_tree.remove(sp)

    # Deep-copy all shapes from source
    src_tree = template.shapes._spTree
    for elem in src_tree:
        sp_tree.append(copy.deepcopy(elem))

    # Copy background
    src_bg = template._element.find(qn('p:cSld'))
    dst_bg = new_slide._element.find(qn('p:cSld'))
    if src_bg is not None and dst_bg is not None:
        src_bg_fill = src_bg.find(qn('p:bg'))
        dst_bg_fill = dst_bg.find(qn('p:bg'))
        if src_bg_fill is not None:
            new_bg = copy.deepcopy(src_bg_fill)
            if dst_bg_fill is not None:
                dst_bg.replace(dst_bg_fill, new_bg)
            else:
                dst_bg.insert(0, new_bg)

    # Reorder: move last slide to insert_at
    sldIdLst = dest_prs.slides._sldIdLst
    last = sldIdLst[-1]
    sldIdLst.remove(last)
    sldIdLst.insert(insert_at, last)

    return new_slide

new_slide = inject_slide(main, src, 0, insert_at=5)

# Add speaker notes
notes_tf = new_slide.notes_slide.notes_text_frame
notes_tf.text = (
    "SCENARIO SETUP SLIDE — spend ~2 minutes here to anchor the audience.\n\n"
    "This is the running example for the entire data modeling section.\n"
    "Three entities from a familiar domain — e-commerce order management.\n\n"
    "Walk through each entity:\n"
    "• Customer: who buys. Has a region (APAC/EMEA/AMER) and a tier (gold/silver/bronze).\n"
    "• Order: what was purchased. Belongs to one customer (1:N relationship). "
    "Has a status (pending/completed/shipped) and a total amount.\n"
    "• Order Item: the line items inside an order. "
    "Each order has 1 or more items (1:N). Stores SKU, quantity, and unit price.\n\n"
    "Point to the relationship arrows:\n"
    "  Customer places many Orders.\n"
    "  Each Order contains many Order Items.\n\n"
    "Then preview what's coming:\n"
    "  01 — SQL: we'll see this as 3 normalized tables with foreign keys and JOINs.\n"
    "  02 — MongoDB: we'll see how this collapses into one document — "
    "Order Items embedded directly inside the Order, Customer referenced by ID.\n"
    "  03 — Every query, index, and transaction demo uses these exact entities "
    "so you can compare SQL vs MongoDB side by side on the same data.\n\n"
    "Say: 'Keep these three entities in your head — every example for the next "
    "40 minutes uses them. By the end you will have seen the same data modeled "
    "two completely different ways and understand why MongoDB makes the choices it does.'"
)

main.save("MongoDB_for_Relational_Developers.pptx")
print(f"✅  Done — total slides: {len(main.slides)}")
os.remove("_scenario_tmp.pptx")
